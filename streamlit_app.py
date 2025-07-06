import streamlit as st
import requests
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader # For PDF parsing
from pptx import Presentation # For PPTX parsing (requires python-pptx)
import io
import os

# For LLM integration:
# import google.generativeai as genai # Commenting out Google Gemini
import openai # Activating OpenAI

# --- API Key Configuration ---
# IMPORTANT: Do NOT hardcode your API key here.
# Use Streamlit Secrets for production deployment.

def extract_text_from_pdf(uploaded_file):
    try:
        reader = PdfReader(io.BytesIO(uploaded_file.read()))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"Error extracting PDF: {e}"

def extract_text_from_pptx(uploaded_file):
    try:
        prs = Presentation(io.BytesIO(uploaded_file.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting PPTX: {e}"

def scrape_website_text(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status() # Raise an HTTPError for bad responses (4xx or 5xx)
        soup = BeautifulSoup(response.content, 'html.parser')
        # Attempt to extract main content, often found in <article>, <main>, or common content divs
        main_content = soup.find('article') or soup.find('main') or soup.find('div', class_='content')
        if main_content:
            return main_content.get_text(separator='\n', strip=True)
        else:
            # Fallback to extracting all visible text if specific content area not found
            texts = soup.stripped_strings
            return "\n".join(texts)
    except requests.exceptions.RequestException as e:
        st.error(f"Error accessing website {url}: {e}")
        return ""
    except Exception as e:
        st.error(f"Error parsing website {url}: {e}")
        return ""

def generate_memo_section_llm(section_name, prompt, document_text, api_key):
    if not api_key:
        return f"Please provide an API key."

    try:
        # OpenRouter setup using OpenAI SDK
        openai.api_key = api_key
        openai.api_base = "https://openrouter.ai/api/v1"

        response = openai.ChatCompletion.create(
            model="openai/gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert investment analyst assistant. Your task is to extract and summarize information from provided documents to create sections of an investment memo. Be concise, factual, and answer the prompt based *only* on the provided text."
                },
                {
                    "role": "user",
                    "content": f"{prompt}\n\nHere is the relevant document text to analyze:\n\n{document_text}"
                }
            ],
            temperature=0.7,
            max_tokens=1000
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        st.error(f"Error generating {section_name}: {e}")
        return f"Could not generate {section_name}. Error: {e}"



# --- Streamlit UI ---
st.set_page_config(layout="wide")
st.title("AI-Powered Investment Memo Generator")
st.markdown("""
Upload a company's pitch deck (PDF/PPTX) and/or provide their website URL to generate an investment memo.
The AI will extract key information, and you'll have the opportunity to review and edit it before finalizing.
""")

# When deployed to Streamlit Cloud, the API key will be loaded from secrets.toml
# Ensure your secret is named "openai_api_key" (or "openrouter_api_key" if you rename it everywhere)
openai_api_key = st.secrets["openai_api_key"] if "openai_api_key" in st.secrets else None

# --- Input Section ---
st.header("1. Input Company Information")
uploaded_pitch_deck = st.file_uploader("Upload Pitch Deck (PDF or PPTX)", type=["pdf", "pptx"])
website_url = st.text_input("Company Website URL (e.g., https://www.example.com)", "")

all_document_text = ""

if uploaded_pitch_deck:
    with st.spinner("Processing pitch deck..."):
        if uploaded_pitch_deck.name.endswith(".pdf"):
            all_document_text += extract_text_from_pdf(uploaded_pitch_deck)
        elif uploaded_pitch_deck.name.endswith(".pptx"):
            all_document_text += extract_text_from_pptx(uploaded_pitch_deck)
    st.success("Pitch deck processed!")

if website_url:
    with st.spinner("Scraping website..."):
        all_document_text += "\n\n" + scrape_website_text(website_url)
    st.success("Website scraped!")

if st.button("Generate Memo"):
    if not uploaded_pitch_deck and not website_url:
        st.warning("Please upload a pitch deck or provide a website URL.")
    elif not all_document_text.strip():
        st.error("Could not extract any meaningful text from the provided sources. Please check the files/URL.")
    elif not openai_api_key: # Checking for the API key in secrets
        st.error("API Key not found in Streamlit secrets. Please configure it for the AI to work.")
    else:
        st.header("2. Review and Edit Investment Memo")

        memo_sections = {
            "Executive Summary": "Generate a concise, plain-language elevator pitch for the company.",
            "Quick Facts": "Extract the company's foundation date, details about any funding rounds (amount, type, investors), and any significant contracts or partnerships mentioned.",
            "Customer Persona": "Describe the ideal customer or target audience for the company's product/service. Who buys it and why?",
            "Problem": "Identify the core problem or pain point that the company's product/service is designed to solve.",
            "Solution": "Describe the actual product, service, or approach the company uses to solve the identified problem. Explain its key features and how it works.",
            "Customer Voice / Expert Opinion": "Extract any customer testimonials, case studies, or expert opinions that validate the company's offering.",
            "Founding Team": "List the key founding team members and their roles, along with any notable relevant experience.",
            "Fundraising and GTM": "Detail the company's fundraising 'ask' (what they are seeking), their go-to-market strategy, and any current traction (e.g., users, revenue, growth metrics).",
            "Key Risk": "Identify and describe the main risks or challenges the company might face (e.g., market competition, regulatory, execution).",
            "Media (Optional)": "Extract any mentions of media coverage, awards, or significant public recognition."
        }

        generated_memo_content = {}
        for section, prompt in memo_sections.items():
            with st.expander(f"Generating: {section}"):
                with st.spinner(f"Extracting {section}..."):
                    generated_content = generate_memo_section_llm(section, prompt, all_document_text, openai_api_key)
                    st.success(f"Extracted {section}!")
            # Create a text area for each section, pre-filled with generated content
            generated_memo_content[section] = st.text_area(
                f"**{section}**",
                value=generated_content,
                height=200,
                key=section # Unique key for each text area
            )

        st.header("3. Final Memo Output")
        final_memo_text = ""
        for section, content in generated_memo_content.items():
            final_memo_text += f"## {section}\n"
            final_memo_text += f"{content}\n\n"

        st.text_area("Full Investment Memo (Editable)", value=final_memo_text, height=600)

        st.download_button(
            label="Download Memo as Markdown",
            data=final_memo_text,
            file_name="investment_memo.md",
            mime="text/markdown"
        )

st.sidebar.header("About This Agent")
st.sidebar.info("""
This AI agent assists in generating initial investment memos by extracting information from pitch decks and websites.
It relies on an LLM (via OpenRouter) for intelligent extraction and provides an editable interface.
""")
st.markdown("""---""")
